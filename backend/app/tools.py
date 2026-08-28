"""Tools the agent can reach for mid-session.

These are real where reality is reachable without a user OAuth consent screen:
- web_search runs live Google Search grounding through Gemini.
- drive_read parses an actual workbook and cites the cells it read.
- slides_export builds a real .pptx you can open in PowerPoint.
- calendar_hold writes a real .ics; gmail_draft writes a real .eml. Both open in
  the native apps; pushing them into a Google account is the one step that needs
  the OAuth consent flow (app/oauth.py), so their summaries say "saved locally".

`execute()` is the single entry point, so adding a tool never touches the loop.
"""
import logging
import re
import uuid
from dataclasses import dataclass
from datetime import datetime, timedelta
from email.message import EmailMessage
from pathlib import Path
from typing import Any, Callable, Optional

from .config import settings
from . import demo_drive
from .images import MEDIA_DIR
from .schemas import SourceFact

log = logging.getLogger(__name__)


@dataclass
class ToolResult:
    ok: bool
    #: One line for the session feed — this is what the group actually reads.
    summary: str
    #: Facts the tool discovered, which become citable on the artifact.
    facts: list[SourceFact]
    #: Anything the caller needs structurally (file urls, draft ids).
    data: dict[str, Any]
    #: True when the real external service was not touched.
    simulated: bool = False


# --------------------------------------------------------------------- tools

def web_search(query: str, uid: str = "") -> ToolResult:
    """Live Google Search grounding: real results, real source domains."""
    from google import genai
    from google.genai import types as gt

    client = genai.Client(
        api_key=settings().google_api_key,
        http_options=gt.HttpOptions(timeout=60_000),
    )
    resp = client.models.generate_content(
        model=settings().gemini_model,
        contents=(
            f"Research: {query}\n"
            "Reply with the 2-3 most useful hard figures, STRICTLY one per line "
            "as `label = value` (short label, concrete value). No other text."
        ),
        config=gt.GenerateContentConfig(
            tools=[gt.Tool(google_search=gt.GoogleSearch())],
            thinking_config=gt.ThinkingConfig(thinking_budget=0),
        ),
    )
    domains: list[str] = []
    gm = resp.candidates[0].grounding_metadata if resp.candidates else None
    if gm and gm.grounding_chunks:
        for c in gm.grounding_chunks:
            if c.web and c.web.title and c.web.title not in domains:
                domains.append(c.web.title)
    facts = []
    for line in (resp.text or "").splitlines():
        m = re.match(r"[-*\s]*([^=|]{3,80})=\s*(.+)", line)
        if m and re.search(r"\d", m.group(2)):
            ref = f"web · {domains[0]}" if domains else "web · google search"
            facts.append(SourceFact(
                fact=m.group(1).strip(" `*"), value=m.group(2).strip(" `*"),
                source_doc="google-search", source_ref=ref,
            ))
    shown = ", ".join(domains[:3]) or "the web"
    return ToolResult(
        True, f"Searched the web — {len(facts)} figures via {shown}",
        facts[:4], {"query": query, "sources": domains[:6]},
    )


def drive_read(file_hint: str, uid: str = "") -> ToolResult:
    """Parse the matching file in the connected drive folder for real."""
    name, facts = demo_drive.read_file(file_hint)
    if not name:
        return ToolResult(False, "No matching file in the drive folder", [], {})
    return ToolResult(
        True, f"Read {name} — pulled {len(facts)} figures",
        facts, {"file": name},
    )


def image_create(query: str, uid: str = "") -> ToolResult:
    """Generate one illustrative image for a page."""
    from . import images
    url, simulated = images.generate(query)
    summary = (
        "Image API unavailable — placed a styled stand-in visual"
        if simulated else "Generated a visual for the page"
    )
    return ToolResult(True, summary, [], {"image_url": url}, simulated=simulated)


def gmail_draft(to: list[str], subject: str, body: str, uid: str = "") -> ToolResult:
    """Write a real .eml draft. Pushing it into Gmail needs the OAuth consent
    flow, so until then it is saved locally and never sent by anyone but you."""
    msg = EmailMessage()
    msg["From"] = "you@deckhand.demo"
    msg["To"] = ", ".join(to) or "the room"
    msg["Subject"] = subject or "Follow-up"
    msg.set_content(body or "")
    name = f"draft-{uuid.uuid4().hex[:8]}.eml"
    (MEDIA_DIR / name).write_bytes(bytes(msg))
    return ToolResult(
        True,
        f"Drafted “{subject}” to {', '.join(to) or 'the room'} — waiting for you to send",
        [], {"url": _url(name), "link_label": "open draft"},
        simulated=True,
    )


def calendar_hold(title: str, attendees: list[str], when: str, uid: str = "") -> ToolResult:
    """Write a real .ics hold that opens straight into Calendar."""
    start = _parse_when(when)
    end = start + timedelta(minutes=30)
    stamp = "%Y%m%dT%H%M%S"
    ics = "\r\n".join([
        "BEGIN:VCALENDAR", "VERSION:2.0", "PRODID:-//Deckhand//EN",
        "BEGIN:VEVENT",
        f"UID:{uuid.uuid4().hex}@deckhand",
        f"DTSTART:{start.strftime(stamp)}",
        f"DTEND:{end.strftime(stamp)}",
        f"SUMMARY:{title or 'Working session'}",
        f"DESCRIPTION:Held by Deckhand for {', '.join(attendees) or 'the room'}",
        "END:VEVENT", "END:VCALENDAR", "",
    ])
    name = f"hold-{uuid.uuid4().hex[:8]}.ics"
    (MEDIA_DIR / name).write_text(ics)
    pretty = start.strftime("%a %b %-d, %-I:%M %p")
    return ToolResult(
        True, f"Held “{title}” — {pretty} with {max(len(attendees), 1)} people",
        [], {"url": _url(name), "link_label": "add to calendar", "when": pretty},
        simulated=True,
    )


def slides_export(title: str, page_count: int = 0, uid: str = "",
                  ctx: Optional[dict] = None) -> ToolResult:
    """Build a real .pptx from the artifact pages. Opens in PowerPoint/Slides."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    pages = (ctx or {}).get("pages", [])
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    ink, mut, brand = RGBColor(0x14, 0x16, 0x1F), RGBColor(0x5A, 0x60, 0x72), RGBColor(0x5B, 0x67, 0xF2)

    def _text(slide, x, y, w, h, s, size, color, bold=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = s
        f = run.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
        return box

    for page in pages:
        slide = prs.slides.add_slide(blank)
        html = page.get("body", "")
        h1 = re.search(r"<h1[^>]*>(.*?)</h1>", html, re.S)
        headline = re.sub(r"<[^>]+>", "", h1.group(1)).strip() if h1 else page.get("label", "")
        para = re.search(r"<p[^>]*>(.*?)</p>", html, re.S)
        support = re.sub(r"<[^>]+>", "", para.group(1)).strip() if para else ""
        img = re.search(r'<img[^>]+src="[^"]*/media/([^"]+)"', html)
        if img and (MEDIA_DIR / img.group(1)).exists():
            try:
                slide.shapes.add_picture(str(MEDIA_DIR / img.group(1)), 0, 0,
                                         prs.slide_width, prs.slide_height)
            except Exception:  # noqa: BLE001 - a bad image must not sink the export
                pass
        light = bool(img)
        _text(slide, 0.9, 0.7, 6, 0.4, page.get("label", "").upper(), 13, brand, bold=True)
        _text(slide, 0.85, 1.5, 11.5, 2.6, headline, 40,
              RGBColor(0xFF, 0xFF, 0xFF) if light else ink, bold=True)
        if support:
            _text(slide, 0.9, 4.2, 10, 1.6, support, 20,
                  RGBColor(0xD8, 0xDA, 0xE6) if light else mut)
    name = f"{re.sub(r'[^a-z0-9]+', '-', (title or 'deck').lower()).strip('-')}-{uuid.uuid4().hex[:6]}.pptx"
    prs.save(MEDIA_DIR / name)
    return ToolResult(
        True, f"Exported {len(pages)} slides to “{title or 'the deck'}” (.pptx)",
        [], {"url": _url(name), "link_label": "download deck"},
    )


def _url(name: str) -> str:
    return f"{settings().public_base_url}/media/{name}"


def _parse_when(when: str) -> datetime:
    """Loose natural parse; a demo must never die on a date format."""
    now = datetime.now().replace(minute=0, second=0, microsecond=0)
    text = (when or "").lower()
    days = ["monday", "tuesday", "wednesday", "thursday", "friday", "saturday", "sunday"]
    target = now + timedelta(days=1)
    for i, day in enumerate(days):
        if day in text:
            ahead = (i - now.weekday()) % 7 or 7
            target = now + timedelta(days=ahead)
            break
    hour = 9
    m = re.search(r"(\d{1,2})(?::(\d{2}))?\s*(am|pm)?", text)
    if m:
        hour = int(m.group(1)) % 12 + (12 if m.group(3) == "pm" else 0)
        if not m.group(3) and int(m.group(1)) < 8:
            hour = int(m.group(1)) + 12
    return target.replace(hour=hour)


# ------------------------------------------------------------------ registry

ToolFn = Callable[..., ToolResult]

REGISTRY: dict[str, ToolFn] = {
    "web_search": web_search,
    "drive_read": drive_read,
    "image_create": image_create,
    "gmail_draft": gmail_draft,
    "calendar_hold": calendar_hold,
    "slides_export": slides_export,
}

#: Given to the conductor so it knows what it is allowed to reach for.
CATALOGUE = """- web_search(query): LIVE Google Search. Use when the group needs a number or
  claim nobody in the room has — a benchmark, market size, competitor fact.
- drive_read(file_hint): read a file from the connected drive folder (revenue
  workbook, cohort data) and pull real figures with their cell references.
- image_create(query): generate one illustrative image for a page. Put the full
  art direction in query (subject, mood, composition) — never words or text in
  the image.
- gmail_draft(to, subject, body): write a follow-up email draft. Never sends.
- calendar_hold(title, attendees, when): hold time for an agreed follow-up.
- slides_export(title, page_count): build the deck as a real .pptx download."""


def execute(name: str, args: dict[str, Any], uid: str = "",
            ctx: Optional[dict] = None) -> Optional[ToolResult]:
    """Run a tool by name. Unknown tools and bad arguments fail soft."""
    fn = REGISTRY.get(name)
    if fn is None:
        log.warning("agent asked for unknown tool %r", name)
        return None
    if name == "slides_export":
        args = {**args, "ctx": ctx}
    try:
        return fn(uid=uid, **args)
    except TypeError as exc:
        log.warning("bad arguments for tool %s: %s", name, exc)
        return None
    except Exception:  # noqa: BLE001
        log.exception("tool %s failed", name)
        return ToolResult(False, f"{name} hit a snag — carrying on without it", [], {})


#: Which flat ToolCall fields each tool actually consumes.
ARG_FIELDS: dict[str, tuple[str, ...]] = {
    "web_search": ("query",),
    "drive_read": ("file_hint",),
    "image_create": ("query",),
    "gmail_draft": ("to", "subject", "body"),
    "calendar_hold": ("title", "attendees", "when"),
    "slides_export": ("title", "page_count"),
}


def args_for(call: Any) -> dict[str, Any]:
    """Pick just the fields the named tool takes off a flat ToolCall."""
    return {field: getattr(call, field) for field in ARG_FIELDS.get(call.name, ())}
